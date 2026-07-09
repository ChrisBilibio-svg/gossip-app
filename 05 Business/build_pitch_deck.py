"""
Fofoca investor deck generator — v1.0
Builds a branded 14-slide investor deck (.pptx).
Run:  python build_pitch_deck.py   ->  FOFOCA_INVESTOR_DECK.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Brand
PINK = RGBColor(0xFF, 0x4D, 0x9D)
DARK = RGBColor(0x0F, 0x14, 0x19)
GREY = RGBColor(0x53, 0x67, 0x71)
LIGHT = RGBColor(0xFF, 0xF3, 0xF8)
TEAL = RGBColor(0x14, 0xB8, 0xA6)
YELLOW = RGBColor(0xFF, 0xD4, 0x3B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(1, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
    return s


def box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def para(tf, text, size, color, bold=False, first=False, align=PP_ALIGN.LEFT,
         space_after=8, bullet=False, italic=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    run = p.add_run(); run.text = text
    f = run.font
    f.size = Pt(size); f.bold = bold; f.italic = italic; f.color.rgb = color
    f.name = "Segoe UI"
    if bullet:
        p.text = ""  # reset; build bullet manually
        run = p.add_run(); run.text = "•  " + text
        run.font.size = Pt(size); run.font.color.rgb = color; run.font.name = "Segoe UI"; run.font.bold = bold
    return p


def kicker(slide, text):
    tf = box(slide, Inches(0.7), Inches(0.5), Inches(11), Inches(0.5))
    para(tf, text.upper(), 13, PINK, bold=True, first=True)


def title(slide, text, size=34):
    tf = box(slide, Inches(0.7), Inches(0.9), Inches(12), Inches(1.1))
    para(tf, text, size, DARK, bold=True, first=True)


def bullets(slide, items, top=2.3, size=18, left=0.75, width=11.8, color=DARK):
    tf = box(slide, Inches(left), Inches(top), Inches(width), Inches(4.4))
    for i, it in enumerate(items):
        para(tf, it, size, color, first=(i == 0), bullet=True, space_after=12)


def footer(slide, n):
    tf = box(slide, Inches(0.7), Inches(7.0), Inches(12), Inches(0.4))
    para(tf, f"Fofoca · confidential · v1.0 · {n}", 9, GREY, first=True)


# ---------------- 1. Title ----------------
s = add_slide(DARK)
tf = box(s, Inches(0.9), Inches(2.6), Inches(11.5), Inches(2))
para(tf, "Fofoca", 72, WHITE, bold=True, first=True)
para(tf, "Brazil's gossip prediction game", 26, PINK, bold=True)
tf2 = box(s, Inches(0.95), Inches(4.7), Inches(11), Inches(1.5))
para(tf2, "Curated rumors  →  the crowd predicts  🍵 TEA / 🧢 CAP  →  credible sources resolve.",
     18, RGBColor(0xCF, 0xD9, 0xDE), first=True)
para(tf2, "A status game for people who are always “right first.”", 16, GREY, italic=True)

# ---------------- 2. Problem ----------------
s = add_slide(); kicker(s, "The problem")
title(s, "Brazilian gossip attention is massive — but fragmented")
bullets(s, [
    "Fans chase rumors across Instagram, X, blogs, TV, and group chats — nothing ties it together.",
    "“Rumor” vs “confirmed fact” is blurry, which erodes trust and fuels misinformation.",
    "Reading gossip is passive. Nothing rewards being early, or being right.",
    "There's no scoreboard for the people whose instincts are consistently correct.",
])
footer(s, "2")

# ---------------- 3. Insight ----------------
s = add_slide(LIGHT); kicker(s, "The insight")
title(s, "The truth filter can become the game")
bullets(s, [
    "Every Speculated rumor is an open prediction the crowd can call.",
    "The crowd doesn't assert facts — it expresses probabilities (TEA = true, CAP = cap/false).",
    "When credible evidence lands, correct early callers earn points, rank, and status.",
    "Skill-weighted scoring rewards contrarian-and-correct, not just loud.",
])
footer(s, "3")

# ---------------- 4. Product ----------------
s = add_slide(); kicker(s, "The product")
title(s, "A daily loop built around open bets")
bullets(s, [
    "1 + 6 curated daily format: one hero “Fofoca do Dia” plus supporting rumors.",
    "Tap 🍵 TEA or 🧢 CAP, lock your call, then the crowd split is revealed.",
    "Article view: headline → short summary → full story → comments.",
    "Resolution payoff: points, leaderboard movement, and bragging rights.",
    "Anonymous-first: play instantly, save your account later.",
])
footer(s, "4")

# ---------------- 5. How it works ----------------
s = add_slide(LIGHT); kicker(s, "How it works")
title(s, "Evidence-first resolution, skill-weighted scoring")
bullets(s, [
    "Curators publish rumors from free sources (AI-summarized drafts, human-approved).",
    "Source credibility is tagged: reliable → Confirmed, social → Speculated.",
    "Rumors resolve when 2+ credible sources confirm/debunk — or by an explicit deadline.",
    "Scoring rewards being early and contrarian-correct; money never touches the truth meter.",
])
footer(s, "5")

# ---------------- 6. Why now ----------------
s = add_slide(); kicker(s, "Why now")
title(s, "Brazil is primed for this mechanic")
bullets(s, [
    "Mobile-native celebrity and fandom culture at enormous scale.",
    "“Palpite” / betting behavior makes prediction intuitive and social.",
    "AI makes summarization and curation leverage cheap — humans keep legal judgment.",
    "No incumbent owns “the gossip scoreboard.”",
])
footer(s, "6")

# ---------------- 7. Market & wedge ----------------
s = add_slide(LIGHT); kicker(s, "Market & wedge")
title(s, "Start narrow, become the default gossip game")
bullets(s, [
    "Beachhead: one dense community — BBB / reality TV, celebrity couples, football transfers, or a creator ecosystem.",
    "Win the daily habit in one wedge before broadening.",
    "Expansion path: broader entertainment → sports rumors → viral news prediction.",
    "Cold-start is the central risk — so we launch where the crowd already gathers.",
])
footer(s, "7")

# ---------------- 8. Business model ----------------
s = add_slide(); kicker(s, "Business model")
title(s, "Don't bet the company on ads")
bullets(s, [
    "Free core game — maximum liquidity and habit. We never charge to predict.",
    "Light ads after retention — secondary line, not the thesis.",
    "Pro (R$ 19.90/mo) — early alerts, advanced stats, status cosmetics. Never pay-to-win.",
    "Gossip Pulse (B2B) — anonymized belief/trend dashboards for media, creators, agencies.",
])
footer(s, "8")

# ---------------- 9. Traction ----------------
s = add_slide(LIGHT); kicker(s, "Traction / build status")
title(s, "v1 is already code-complete enough to test")
bullets(s, [
    "Expo + Supabase app: feed, betting, evidence-first resolution, leaderboard, comments, admin.",
    "Content automation drafts curator-ready rumors from free sources (cheap AI summarization).",
    "Source-credibility tagging, prediction deadlines, and keyword search shipped.",
    "The next gap is not code — it's beachhead launch, analytics, and legal polish.",
])
footer(s, "9")

# ---------------- 10. Financials: scenarios ----------------
s = add_slide(); kicker(s, "Financial model")
title(s, "Lean base case, real upside (36-month plan)")
tf = box(s, Inches(0.75), Inches(2.2), Inches(11.8), Inches(3.2))
rows = [
    ("Scenario", "M36 MAU", "M36 Revenue/mo", "M36 ARR"),
    ("Conservative", "6,142", "R$ 17,198", "R$ 0.21M"),
    ("Base", "133,081", "R$ 111,601", "R$ 1.34M"),
    ("Upside", "889,798", "R$ 702,589", "R$ 8.43M"),
]
for i, row in enumerate(rows):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(10)
    run = p.add_run(); run.text = "    ".join(f"{c:<16}" for c in row)
    run.font.name = "Consolas"; run.font.size = Pt(16)
    run.font.bold = (i == 0)
    run.font.color.rgb = PINK if i == 0 else DARK
para(box(s, Inches(0.75), Inches(5.4), Inches(11.8), Inches(1)),
     "Month 1 = 1,000-MAU private beta, not a viral launch. Revenue stack: ads + Pro + B2B Pulse.",
     14, GREY, first=True, italic=True)
footer(s, "10")

# ---------------- 11. Two paths ----------------
s = add_slide(LIGHT); kicker(s, "Two strategic paths")
title(s, "Bootstrapped vs Funded")
tf = box(s, Inches(0.75), Inches(2.2), Inches(11.8), Inches(3.6))
rows = [
    ("", "Bootstrapped", "Funded"),
    ("M36 MAU", "133,081", "889,798"),
    ("M36 revenue/mo", "R$ 111.6k", "R$ 702.6k"),
    ("M36 net/mo", "+R$ 36.5k", "+R$ 422.1k"),
    ("Capital raised", "R$ 0 (needs ~R$752k runway)", "R$ 9.5M (Seed + Series A)"),
    ("Founder ownership", "~90%", "~57.6%"),
    ("Founder value @6x ARR", "R$ 7.2M (US$ 1.36M)", "R$ 29.1M (US$ 5.5M)"),
]
for i, row in enumerate(rows):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(8)
    run = p.add_run(); run.text = f"{row[0]:<24}{row[1]:<30}{row[2]:<26}"
    run.font.name = "Consolas"; run.font.size = Pt(14)
    run.font.bold = (i == 0)
    run.font.color.rgb = PINK if i == 0 else DARK
footer(s, "11")

# ---------------- 12. The ask ----------------
s = add_slide(DARK); kicker(s, "The ask")
tf = box(s, Inches(0.7), Inches(0.9), Inches(12), Inches(1.1))
para(tf, "Raising a Seed round to win the beachhead", 32, WHITE, bold=True, first=True)
bullets(s, [
    "Seed target: ~R$ 1.5M (≈ US$ 280k) to fund the validation + early-growth ramp.",
    "Use of funds: curation + community, one growth hire, mobile/backend, legal/trust.",
    "Goal: prove D7 prediction retention and bets/DAU in one dense wedge.",
    "Milestone to Series A: durable daily habit + first Gossip Pulse B2B pilots.",
], top=2.4, color=RGBColor(0xE6, 0xEC, 0xF0))
footer(s, "12")

# ---------------- 13. Team ----------------
s = add_slide(); kicker(s, "Team")
title(s, "Founder-led, lean, AI-leveraged")
bullets(s, [
    "Chris — founder / product & business; building the managed-AI workflow behind the app.",
    "Co-dev (Pedro) — engineering partner; shipped the evidence-first resolution + search.",
    "AI-leveraged ops: content summarization and curation drafting run cheaply, humans keep judgment.",
    "Hiring plan ramps curation, growth, engineering, community, and BD as the model scales.",
])
footer(s, "13")

# ---------------- 14. Close ----------------
s = add_slide(PINK)
tf = box(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(2.5))
para(tf, "The gossip scoreboard for Brazil.", 40, WHITE, bold=True, first=True)
para(tf, "Free to play. Status to earn. Trust by design.", 22, WHITE)
para(tf, "contato@fofoca.app", 18, DARK, bold=True, space_after=2)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "FOFOCA_INVESTOR_DECK.pptx")
prs.save(out)
print("Saved:", out, "·", len(prs.slides._sldIdLst), "slides")
